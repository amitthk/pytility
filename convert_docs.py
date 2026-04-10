#!/usr/bin/env python3
"""
convert_docs.py
───────────────
Batch-convert DOCX, PPTX, and PDF files to Markdown / plain-text.

  DOCX  →  Markdown  (headings, bullets, tables, embedded images)
  PPTX  →  Markdown  (slide titles, text, images)
  PDF   →  Markdown or plain text  (page-by-page via PyMuPDF)

Images are saved to <output_dir>/images/ and referenced from the
generated Markdown with relative paths.

Usage
─────
  python convert_docs.py                          # current dir, all types
  python convert_docs.py ./docs -o ./out
  python convert_docs.py report.docx -o ./out
  python convert_docs.py ./docs --types docx,pptx
  python convert_docs.py ./docs --types pdf --pdf-format markdown
  python convert_docs.py ./docs -o ./out --recursive --skip-existing -v

Dependencies
────────────
  pip install python-docx python-pptx pymupdf
"""

import argparse
import re
import shutil
import sys
import zipfile
from pathlib import Path


# ─────────────────────────────────────────────────────────────────────────────
# Image helpers
# ─────────────────────────────────────────────────────────────────────────────

SUPPORTED_IMG_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".emf", ".wmf"}


def safe_stem(path: Path) -> str:
    return path.stem.replace(" ", "_")


def extract_zip_images(archive_path: Path, media_prefix: str,
                       img_dir: Path, file_prefix: str) -> dict[str, str]:
    """
    Extract images from a zip-based Office file.
    Returns {internal_name: saved_filename}.
    """
    saved = {}
    with zipfile.ZipFile(archive_path) as z:
        media_files = [n for n in z.namelist() if n.startswith(media_prefix)]
        for mf in media_files:
            if Path(mf).suffix.lower() not in SUPPORTED_IMG_EXTS:
                continue
            out_name = f"{file_prefix}_{Path(mf).name}"
            out_path = img_dir / out_name
            with z.open(mf) as src, open(out_path, "wb") as dst:
                shutil.copyfileobj(src, dst)
            saved[Path(mf).name] = out_name
    return saved


# ─────────────────────────────────────────────────────────────────────────────
# PPTX → Markdown
# ─────────────────────────────────────────────────────────────────────────────

def convert_pptx(pptx_path: Path, out_dir: Path, img_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Pt

    prefix = safe_stem(pptx_path)
    extract_zip_images(pptx_path, "ppt/media/", img_dir, prefix)

    prs = Presentation(pptx_path)
    lines = [f"# {pptx_path.stem}\n"]

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = ""
        text_blocks = []
        img_refs = []

        for shape in slide.shapes:
            # ── text ──────────────────────────────────────────────────────────
            if shape.has_text_frame and shape.shape_type != 13:
                tf = shape.text_frame
                for para_idx, para in enumerate(tf.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    is_heading = False
                    try:
                        if shape.is_title:
                            is_heading = True
                    except Exception:
                        pass
                    if not is_heading and para.runs:
                        try:
                            sz = para.runs[0].font.size
                            if sz and sz >= Pt(18):
                                is_heading = True
                        except Exception:
                            pass

                    if is_heading and para_idx == 0 and not slide_title:
                        slide_title = text
                    elif is_heading:
                        text_blocks.append(f"### {text}")
                    else:
                        level = para.level or 0
                        bullet = "  " * level + "- "
                        text_blocks.append(f"{bullet}{text}")

            # ── images ────────────────────────────────────────────────────────
            if shape.shape_type == 13:
                try:
                    blob = shape.image.blob
                    ext  = shape.image.ext
                    img_name = f"{prefix}_slide{slide_num:03d}_img{len(img_refs) + 1}.{ext}"
                    with open(img_dir / img_name, "wb") as f:
                        f.write(blob)
                    img_refs.append(img_name)
                except Exception as e:
                    img_refs.append(f"[image extraction error: {e}]")

        lines.append(f"\n## Slide {slide_num}{': ' + slide_title if slide_title else ''}\n")
        lines.extend(text_blocks)
        for img_name in img_refs:
            if img_name.startswith("["):
                lines.append(f"\n{img_name}\n")
            else:
                rel = _rel_img_path(out_dir, img_dir, img_name)
                lines.append(f"\n![Slide {slide_num} image]({rel})\n")

    out_path = out_dir / (prefix + ".md")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# DOCX → Markdown
# ─────────────────────────────────────────────────────────────────────────────

def convert_docx(docx_path: Path, out_dir: Path, img_dir: Path) -> Path:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    prefix = safe_stem(docx_path)
    extract_zip_images(docx_path, "word/media/", img_dir, prefix)

    doc = Document(docx_path)
    lines = [f"# {docx_path.stem}\n"]
    img_counter = [0]

    def para_to_md(para) -> str:
        style = para.style.name if para.style else ""
        text  = para.text.strip()

        imgs = []
        for run in para.runs:
            for drawing in run._element.findall(".//" + qn("a:blip")):
                rId = drawing.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                if rId:
                    try:
                        part = doc.part.related_parts[rId]
                        ext  = Path(part.partname).suffix.lstrip(".")
                        img_counter[0] += 1
                        img_name = f"{prefix}_img{img_counter[0]:03d}.{ext}"
                        with open(img_dir / img_name, "wb") as f:
                            f.write(part.blob)
                        rel = _rel_img_path(out_dir, img_dir, img_name)
                        imgs.append(f"![image]({rel})")
                    except Exception as e:
                        imgs.append(f"[image: {e}]")

        parts = []
        if text:
            if re.match(r"Heading [1-6]", style):
                level = int(style.split()[-1])
                parts.append("#" * level + " " + text)
            elif style == "Title":
                parts.append("# " + text)
            elif re.match(r"List Bullet", style):
                depth = style[-1] if style[-1].isdigit() else "1"
                parts.append("  " * (int(depth) - 1) + "- " + text)
            elif re.match(r"List Number", style):
                depth = style[-1] if style[-1].isdigit() else "1"
                parts.append("  " * (int(depth) - 1) + "1. " + text)
            else:
                parts.append(text)
        parts.extend(imgs)
        return "\n".join(parts)

    def table_to_md(table) -> str:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("|" + "|".join(["---"] * len(cells)) + "|")
        return "\n".join(rows)

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
        if tag == "p":
            md = para_to_md(Paragraph(block, doc))
            if md:
                lines.append(md)
        elif tag == "tbl":
            lines.append("\n" + table_to_md(Table(block, doc)) + "\n")

    out_path = out_dir / (prefix + ".md")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# PDF → Markdown or plain text
# ─────────────────────────────────────────────────────────────────────────────

def convert_pdf(pdf_path: Path, out_dir: Path, img_dir: Path,
                fmt: str = "text", page_markers: bool = True) -> Path:
    """
    fmt = "text"     → plain text with optional page markers  → .txt
    fmt = "markdown" → markdown with page headings            → .md
    """
    import fitz  # PyMuPDF

    prefix = safe_stem(pdf_path)
    doc    = fitz.open(str(pdf_path))
    pages  = []
    img_counter = 0

    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text").strip()

        # ── extract embedded images when writing markdown ──────────────────
        if fmt == "markdown":
            img_refs = []
            for img_index, img_info in enumerate(page.get_images(full=True)):
                xref = img_info[0]
                try:
                    base_img  = doc.extract_image(xref)
                    img_bytes = base_img["image"]
                    img_ext   = base_img["ext"]
                    img_counter += 1
                    img_name  = f"{prefix}_p{page_num:03d}_img{img_counter}.{img_ext}"
                    with open(img_dir / img_name, "wb") as f:
                        f.write(img_bytes)
                    rel = _rel_img_path(out_dir, img_dir, img_name)
                    img_refs.append(f"![Page {page_num} image {img_index + 1}]({rel})")
                except Exception:
                    pass

            if not text and not img_refs:
                continue

            block_parts = []
            if page_markers:
                block_parts.append(f"## Page {page_num}")
            if text:
                block_parts.append(text)
            block_parts.extend(img_refs)
            pages.append("\n\n".join(block_parts))
        else:
            if not text:
                continue
            header = f"--- Page {page_num} ---\n" if page_markers else ""
            pages.append(header + text)

    doc.close()

    ext      = ".md" if fmt == "markdown" else ".txt"
    out_path = out_dir / (prefix + ext)
    separator = "\n\n---\n\n" if fmt == "markdown" else "\n\n"
    content  = separator.join(pages)
    if fmt == "markdown":
        content = f"# {pdf_path.stem}\n\n" + content
    out_path.write_text(content, encoding="utf-8")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ─────────────────────────────────────────────────────────────────────────────

def convert_file(path: Path, out_dir: Path, img_dir: Path,
                 pdf_format: str, page_markers: bool,
                 verbose: bool) -> tuple[bool, str]:
    """
    Convert a single file. Returns (success, message).
    """
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            out = convert_docx(path, out_dir, img_dir)
        elif ext == ".pptx":
            out = convert_pptx(path, out_dir, img_dir)
        elif ext == ".pdf":
            out = convert_pdf(path, out_dir, img_dir, pdf_format, page_markers)
        else:
            return False, f"unsupported type: {ext}"
        kb = out.stat().st_size // 1024
        msg = f"→ {out.name}  ({kb} KB)"
        return True, msg
    except ImportError as e:
        return False, f"missing dependency — {e}"
    except Exception as e:
        return False, f"error — {e}"


def _rel_img_path(out_dir: Path, img_dir: Path, img_name: str) -> str:
    """Return a relative path from out_dir to the image file."""
    try:
        return str((img_dir / img_name).relative_to(out_dir)).replace("\\", "/")
    except ValueError:
        return str(img_dir / img_name)


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

SUPPORTED_TYPES = {
    "docx": [".docx"],
    "pptx": [".pptx"],
    "pdf":  [".pdf"],
}


def collect_files(input_path: Path, types: list[str], recursive: bool) -> list[Path]:
    extensions = set()
    for t in types:
        extensions.update(SUPPORTED_TYPES[t])

    if input_path.is_file():
        return [input_path] if input_path.suffix.lower() in extensions else []

    glob = "**/*" if recursive else "*"
    return sorted(
        p for p in input_path.glob(glob)
        if p.is_file() and p.suffix.lower() in extensions
    )


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Convert DOCX / PPTX / PDF to Markdown or plain text",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument(
        "input",
        nargs="?",
        default=".",
        help="Input file or directory (default: current directory)",
    )
    p.add_argument(
        "--out", "-o",
        default=None,
        metavar="DIR",
        help="Output directory (default: same as input)",
    )
    p.add_argument(
        "--img-dir",
        default=None,
        metavar="DIR",
        help="Directory for extracted images (default: <out>/images)",
    )
    p.add_argument(
        "--types",
        default="all",
        metavar="TYPE[,TYPE]",
        help="File types to process: docx, pptx, pdf, or all (default: all)",
    )
    p.add_argument(
        "--pdf-format",
        choices=["text", "markdown"],
        default="text",
        help="Output format for PDF files (default: text)",
    )
    p.add_argument(
        "--no-page-markers",
        action="store_true",
        help="Omit page markers in PDF output",
    )
    p.add_argument(
        "--recursive", "-r",
        action="store_true",
        help="Recurse into subdirectories when input is a directory",
    )
    p.add_argument(
        "--skip-existing",
        action="store_true",
        help="Skip files whose output already exists",
    )
    p.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Print per-file detail",
    )
    return p.parse_args()


def main() -> int:
    args = parse_args()

    # ── resolve paths ─────────────────────────────────────────────────────────
    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"[error] input not found: {input_path}", file=sys.stderr)
        return 1

    out_dir = Path(args.out).resolve() if args.out else (
        input_path if input_path.is_dir() else input_path.parent
    )
    out_dir.mkdir(parents=True, exist_ok=True)

    img_dir = Path(args.img_dir).resolve() if args.img_dir else out_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)

    # ── resolve types ─────────────────────────────────────────────────────────
    raw_types = args.types.lower().replace(" ", "")
    if raw_types == "all":
        active_types = list(SUPPORTED_TYPES.keys())
    else:
        active_types = []
        for t in raw_types.split(","):
            if t not in SUPPORTED_TYPES:
                print(f"[error] unknown type '{t}'. Choose from: docx, pptx, pdf, all",
                      file=sys.stderr)
                return 1
            active_types.append(t)

    # ── collect files ─────────────────────────────────────────────────────────
    files = collect_files(input_path, active_types, args.recursive)
    if not files:
        print(f"[warn] no matching files found in: {input_path}")
        return 0

    # ── print header ──────────────────────────────────────────────────────────
    print(f"\nDocument Converter")
    print(f"  input   : {input_path}")
    print(f"  output  : {out_dir}")
    print(f"  images  : {img_dir}")
    print(f"  types   : {', '.join(active_types)}")
    print(f"  pdf fmt : {args.pdf_format}")
    print(f"  files   : {len(files)}")
    print()

    # ── convert ───────────────────────────────────────────────────────────────
    ok_count = fail_count = skip_count = 0

    for file_path in files:
        ext = file_path.suffix.lower().lstrip(".")
        expected_ext = ".md" if ext in ("docx", "pptx") or args.pdf_format == "markdown" else ".txt"
        out_file = out_dir / (safe_stem(file_path) + expected_ext)

        label = file_path.name
        if input_path.is_dir():
            try:
                label = str(file_path.relative_to(input_path))
            except ValueError:
                pass

        if args.skip_existing and out_file.exists():
            print(f"  [skip]  {label}")
            skip_count += 1
            continue

        print(f"  {label}")
        ok, msg = convert_file(
            file_path, out_dir, img_dir,
            pdf_format   = args.pdf_format,
            page_markers = not args.no_page_markers,
            verbose      = args.verbose,
        )
        status = "✓" if ok else "✗"
        print(f"    {status}  {msg}")
        if ok:
            ok_count += 1
        else:
            fail_count += 1

    # ── summary ───────────────────────────────────────────────────────────────
    imgs = list(img_dir.iterdir()) if img_dir.exists() else []
    print()
    print(f"Done — {ok_count} converted, {fail_count} failed, {skip_count} skipped")
    print(f"Images extracted : {len(imgs)}")
    print(f"Output directory : {out_dir}")
    return 0 if fail_count == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
